"""
Reusable slide-deck builder matching the "MLOps MENA" reference deck style:
white background, Bricolage Grotesque headings, Nunito Sans body, IBM Plex
Mono for eyebrows/code/page numbers, VS-Code-dark-style code blocks.

Install the fonts locally for pixel-perfect rendering (all free / Google
Fonts): "Bricolage Grotesque", "Nunito Sans", "IBM Plex Mono". Without them,
PowerPoint substitutes a default font but the layout still holds.

Usage (see Day 01/Slides_01/build.py for a full example):

    from deckkit import Deck

    deck = Deck()
    deck.title_slide("SQL FOUNDATIONS · DAY 1 OF 9", "SQL Foundations",
                      "SELECT · LIMIT/OFFSET · DISTINCT · ORDER BY · COUNT",
                      "world database · ~45 min")
    deck.agenda_slide("DAY 01 · AGENDA", "What we cover today", "Six topics",
                       [("01", "Exploring the Server", "SHOW/DESCRIBE"), ...])
    deck.save("Day 01 - SQL Foundations.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette --

INK = RGBColor(0x14, 0x14, 0x14)
BLUE = RGBColor(0x2F, 0x63, 0xE8)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
GRAY_MUTED = RGBColor(0x9C, 0xA3, 0xAF)
CARD_BORDER = RGBColor(0xE5, 0xE7, 0xEB)
PANEL_BG = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)
CODE_DEFAULT = RGBColor(0xD4, 0xD4, 0xD4)
CODE_KEYWORD = RGBColor(0x56, 0x9C, 0xD6)
CODE_FUNC = RGBColor(0xDC, 0xDC, 0xAA)
CODE_STRING = RGBColor(0xCE, 0x91, 0x78)
CODE_COMMENT = RGBColor(0x6A, 0x99, 0x55)
CODE_NUMBER = RGBColor(0xB5, 0xCE, 0xA8)

TRAFFIC = [RGBColor(0xFF, 0x5F, 0x56), RGBColor(0xFF, 0xBD, 0x2E), RGBColor(0x27, 0xC9, 0x3F)]

TITLE_FONT = "Bricolage Grotesque"
BODY_FONT = "Nunito Sans"
MONO_FONT = "IBM Plex Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)

# SQL keyword buckets for the tiny auto-tokenizer used by code_slide()
_KEYWORDS = {
    "SELECT", "FROM", "WHERE", "AND", "OR", "NOT", "ORDER", "BY", "GROUP",
    "HAVING", "LIMIT", "OFFSET", "AS", "ASC", "DESC", "DISTINCT", "IN",
    "BETWEEN", "LIKE", "IS", "NULL", "JOIN", "INNER", "LEFT", "RIGHT",
    "OUTER", "ON", "UNION", "CASE", "WHEN", "THEN", "ELSE", "END", "USE",
    "SHOW", "DESCRIBE", "DATABASES", "SCHEMAS", "TABLES",
}
_FUNCS = {
    "COUNT", "SUM", "AVG", "MIN", "MAX", "ROUND", "CONCAT", "SUBSTRING",
    "TRIM", "LTRIM", "RTRIM", "LEFT", "RIGHT", "REPLACE", "LOCATE",
    "COALESCE", "IFNULL", "DATE_FORMAT", "YEAR", "MONTH", "DAY", "RANK",
    "ROW_NUMBER", "OVER", "PARTITION",
}


def _set_no_autofit(tf):
    el = tf._txBody
    bodyPr = el.find(qn('a:bodyPr'))
    for tag in ('a:normAutofit', 'a:spAutoFit'):
        existing = bodyPr.find(qn(tag))
        if existing is not None:
            bodyPr.remove(existing)
    bodyPr.append(bodyPr.makeelement(qn('a:noAutofit'), {}))


def _tokenize_sql_line(line):
    """Split a line of SQL into (text, color) runs: keywords/functions/strings/comments colored."""
    if line.strip().startswith("--"):
        return [(line, CODE_COMMENT)]
    tokens = []
    buf = ""
    i = 0
    n = len(line)
    while i < n:
        ch = line[i]
        if ch == "'":
            if buf:
                tokens.extend(_word_tokens(buf))
                buf = ""
            j = i + 1
            while j < n and line[j] != "'":
                j += 1
            j = min(j + 1, n)
            tokens.append((line[i:j], CODE_STRING))
            i = j
            continue
        if ch == "-" and i + 1 < n and line[i + 1] == "-":
            if buf:
                tokens.extend(_word_tokens(buf))
                buf = ""
            tokens.append((line[i:], CODE_COMMENT))
            i = n
            continue
        buf += ch
        i += 1
    if buf:
        tokens.extend(_word_tokens(buf))
    return tokens


def _word_tokens(buf):
    import re
    out = []
    for piece in re.split(r'(\W+)', buf):
        if not piece:
            continue
        up = piece.upper()
        if up in _KEYWORDS:
            out.append((piece, CODE_KEYWORD))
        elif up in _FUNCS:
            out.append((piece, CODE_FUNC))
        elif piece.isdigit():
            out.append((piece, CODE_NUMBER))
        else:
            out.append((piece, CODE_DEFAULT))
    return out


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._blank = self.prs.slide_layouts[6]
        self._n = 0

    # -- low-level helpers --------------------------------------------------

    def _slide(self):
        s = self.prs.slides.add_slide(self._blank)
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        rect.fill.solid()
        rect.fill.fore_color.rgb = WHITE
        rect.line.fill.background()
        rect.shadow.inherit = False
        s.shapes._spTree.remove(rect._element)
        s.shapes._spTree.insert(2, rect._element)
        self._n += 1
        return s

    def _textbox(self, slide, left, top, width, height):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        _set_no_autofit(tf)
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        return tb, tf

    def _run(self, para, text, font=BODY_FONT, size=18, color=INK, bold=False, italic=False):
        r = para.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
        return r

    def _header(self, slide, eyebrow, title, subtitle=None, title_size=40):
        _, tf = self._textbox(slide, MARGIN, Inches(0.5), Inches(10.5), Inches(0.4))
        p = tf.paragraphs[0]
        self._run(p, eyebrow.upper(), font=MONO_FONT, size=13, color=BLUE, bold=True)

        _, tf = self._textbox(slide, MARGIN, Inches(0.85), Inches(11.5), Inches(1.1))
        p = tf.paragraphs[0]
        self._run(p, title, font=TITLE_FONT, size=title_size, color=INK, bold=True)

        y = Inches(1.75)
        if subtitle:
            _, tf = self._textbox(slide, MARGIN, y, Inches(11.5), Inches(0.5))
            p = tf.paragraphs[0]
            self._run(p, subtitle, font=BODY_FONT, size=17, color=GRAY)

    def _page_number(self, slide):
        _, tf = self._textbox(slide, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.3))
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        self._run(tf.paragraphs[0], f"{self._n:02d} / {{TOTAL}}", font=MONO_FONT, size=11, color=GRAY_MUTED)

    def _rounded_rect(self, slide, left, top, width, height, fill=None, line_color=None, radius=0.08):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        if line_color is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line_color
            shp.line.width = Pt(1)
        shp.shadow.inherit = False
        return shp

    # -- diagram primitives -----------------------------------------------

    def _entity_box(self, slide, x, y, w, h, title, sub=None, fill=WHITE, border=BLUE):
        self._rounded_rect(slide, x, y, w, h, fill=fill, line_color=border, radius=0.12)
        if sub:
            _, tf = self._textbox(slide, x + Inches(0.08), y + Inches(0.12), w - Inches(0.16), Inches(0.28))
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._run(tf.paragraphs[0], title, font=BODY_FONT, size=14, color=INK, bold=True)
            _, tf = self._textbox(slide, x + Inches(0.08), y + Inches(0.42), w - Inches(0.16), h - Inches(0.5))
            tf.word_wrap = True
            first = True
            for line in sub.split("\n"):
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.alignment = PP_ALIGN.CENTER
                self._run(p, line, font=MONO_FONT, size=9.5, color=GRAY)
        else:
            _, tf = self._textbox(slide, x + Inches(0.05), y, w - Inches(0.1), h)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._run(tf.paragraphs[0], title, font=BODY_FONT, size=14, color=INK, bold=True)

    def _connector(self, slide, x1, y1, x2, y2, color=GRAY_MUTED, width=1.25):
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)))
        ln.line.color.rgb = color
        ln.line.width = Pt(width)
        ln.shadow.inherit = False
        return ln

    def _tag(self, slide, cx, cy, text, color=BLUE, text_color=WHITE, diameter=0.34):
        d = Inches(diameter)
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(cx - d / 2)), Emu(int(cy - d / 2)), d, d)
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        shp.shadow.inherit = False
        tf = shp.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._run(p, text, font=MONO_FONT, size=11, color=text_color, bold=True)

    def _caption(self, slide, x, y, w, text, size=13, color=GRAY, align=PP_ALIGN.CENTER, bold=False):
        _, tf = self._textbox(slide, x, y, w, Inches(0.4))
        tf.paragraphs[0].alignment = align
        self._run(tf.paragraphs[0], text, font=BODY_FONT, size=size, color=color, bold=bold)

    def _table_grid(self, slide, x, y, w, h, rows=4, cols=3, header_fill=BLUE):
        """A small grid of cells, header row filled blue — the visual shorthand for 'relational table'."""
        row_h = h / rows
        col_w = w / cols
        for r in range(rows):
            for c in range(cols):
                cx = x + col_w * c
                cy = y + row_h * r
                fill = header_fill if r == 0 else WHITE
                shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(cx)), Emu(int(cy)), Emu(int(col_w)), Emu(int(row_h)))
                shp.fill.solid()
                shp.fill.fore_color.rgb = fill
                shp.line.color.rgb = CARD_BORDER
                shp.line.width = Pt(1)
                shp.shadow.inherit = False
                if r > 0:
                    tf = shp.text_frame
                    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(2)
                    p = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    dots = "•" * (3 - c % 2) if (r + c) % 2 == 0 else "—" * (2 + c % 2)
                    self._run(p, dots, font=MONO_FONT, size=9, color=GRAY_MUTED)

    def _document_icon(self, slide, x, y, w, h):
        """A rounded box with brace + key: value lines — the visual shorthand for 'flexible document'."""
        self._rounded_rect(slide, x, y, w, h, fill=WHITE, line_color=CARD_BORDER, radius=0.1)
        lines = ["{", '  "id": 101,', '  "tags": [..],', '  "meta": {..}', "}"]
        _, tf = self._textbox(slide, x + Inches(0.22), y + Inches(0.16), w - Inches(0.44), h - Inches(0.3))
        tf.word_wrap = False
        first = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            color = CODE_FUNC if i in (0, 4) else CODE_STRING
            self._run(p, line, font=MONO_FONT, size=11.5, color=color)

    def _oval(self, slide, cx, cy, w, h, text, fill=WHITE, border=CARD_BORDER):
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(cx - w / 2)), Emu(int(cy - h / 2)), Emu(int(w)), Emu(int(h)))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = border
        shp.line.width = Pt(1)
        shp.shadow.inherit = False
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._run(p, text, font=BODY_FONT, size=11, color=INK)
        return shp

    def _diamond(self, slide, cx, cy, w, h, text, fill=PANEL_BG, border=BLUE):
        shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Emu(int(cx - w / 2)), Emu(int(cy - h / 2)), Emu(int(w)), Emu(int(h)))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = border
        shp.line.width = Pt(1.25)
        shp.shadow.inherit = False
        tf = shp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._run(p, text, font=BODY_FONT, size=11, color=INK, bold=True)
        return shp

    # -- slide types ----------------------------------------------------

    def title_slide(self, eyebrow, title, subtitle, footer=None, logo_mark=None):
        s = self._slide()
        cy = Inches(3.1)
        if logo_mark:
            _, tf = self._textbox(s, Inches(0), cy - Inches(0.9), SLIDE_W, Inches(0.6))
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._run(tf.paragraphs[0], logo_mark, font=TITLE_FONT, size=30, color=INK, bold=True)

        _, tf = self._textbox(s, Inches(0), cy, SLIDE_W, Inches(0.4))
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        self._run(tf.paragraphs[0], eyebrow.upper(), font=MONO_FONT, size=15, color=BLUE, bold=True)

        _, tf = self._textbox(s, Inches(0.5), cy + Inches(0.45), Inches(12.333), Inches(1.3))
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        self._run(tf.paragraphs[0], title, font=TITLE_FONT, size=54, color=INK, bold=True)

        _, tf = self._textbox(s, Inches(0.5), cy + Inches(1.65), Inches(12.333), Inches(0.5))
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        self._run(tf.paragraphs[0], subtitle, font=BODY_FONT, size=19, color=GRAY)

        if footer:
            _, tf = self._textbox(s, Inches(0.5), cy + Inches(2.35), Inches(12.333), Inches(0.4))
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._run(tf.paragraphs[0], footer, font=BODY_FONT, size=14, color=GRAY_MUTED)
        return s

    def agenda_slide(self, eyebrow, title, subtitle, items):
        """items: list of (number_str, heading, description)"""
        s = self._slide()
        self._header(s, eyebrow, title, subtitle)
        top = Inches(2.5)
        row_h = Inches(0.52)
        for i, (num, heading, desc) in enumerate(items):
            y = top + row_h * i
            line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, MARGIN, y, Inches(12.733), y)
            line.line.color.rgb = CARD_BORDER
            line.line.width = Pt(0.75)

            _, tf = self._textbox(s, MARGIN, y, Inches(0.6), row_h)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            self._run(tf.paragraphs[0], num, font=MONO_FONT, size=14, color=BLUE, bold=True)

            _, tf = self._textbox(s, Inches(1.3), y, Inches(4.3), row_h)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            self._run(tf.paragraphs[0], heading, font=BODY_FONT, size=16, color=INK, bold=True)

            _, tf = self._textbox(s, Inches(5.8), y, Inches(6.5), row_h)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            self._run(tf.paragraphs[0], desc, font=BODY_FONT, size=14, color=GRAY)

        y = top + row_h * len(items)
        line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, MARGIN, y, Inches(12.733), y)
        line.line.color.rgb = CARD_BORDER
        line.line.width = Pt(0.75)
        self._page_number(s)
        return s

    def grid_slide(self, eyebrow, title, subtitle, cards, columns=2, accent_numbers=False):
        """cards: list of (heading, desc) laid out in a grid (default 2 columns).
        accent_numbers: if True, prefixes each card with a blue 01/02/... index."""
        s = self._slide()
        self._header(s, eyebrow, title, subtitle)

        top = Inches(2.7)
        gap = Inches(0.25)
        n = len(cards)
        rows = (n + columns - 1) // columns
        card_w = (Inches(12.733) - MARGIN - gap * (columns - 1)) / columns
        card_h = min(Inches(1.35), (Inches(7.0) - top - gap * (rows - 1)) / rows)

        for i, (heading, desc) in enumerate(cards):
            r, c = divmod(i, columns)
            x = MARGIN + (card_w + gap) * c
            y = top + (card_h + gap) * r
            self._rounded_rect(s, x, y, card_w, card_h, fill=WHITE, line_color=CARD_BORDER)
            pad = Inches(0.28)
            _, tf = self._textbox(s, x + pad, y + Inches(0.16), card_w - pad * 2, Inches(0.4))
            head_text = f"{i + 1:02d}   {heading}" if accent_numbers else heading
            if accent_numbers:
                p = tf.paragraphs[0]
                self._run(p, f"{i + 1:02d}  ", font=MONO_FONT, size=15, color=BLUE, bold=True)
                self._run(p, heading, font=BODY_FONT, size=17, color=INK, bold=True)
            else:
                self._run(tf.paragraphs[0], heading, font=BODY_FONT, size=17, color=INK, bold=True)
            _, tf = self._textbox(s, x + pad, y + Inches(0.56), card_w - pad * 2, card_h - Inches(0.7))
            tf.word_wrap = True
            self._run(tf.paragraphs[0], desc, font=BODY_FONT, size=13.5, color=GRAY)

        self._page_number(s)
        return s

    def database_types_slide(self, eyebrow, title, subtitle, relational, nosql):
        """relational / nosql: dicts with 'heading', 'desc', 'examples' (str)."""
        s = self._slide()
        self._header(s, eyebrow, title, subtitle)

        col_w = Inches(5.7)
        gap = Inches(0.35)
        top = Inches(2.65)
        for i, (kind, spec) in enumerate([("relational", relational), ("nosql", nosql)]):
            x = MARGIN + (col_w + gap) * i
            self._rounded_rect(s, x, top, col_w, Inches(4.15), fill=WHITE, line_color=CARD_BORDER)

            _, tf = self._textbox(s, x + Inches(0.3), top + Inches(0.22), col_w - Inches(0.6), Inches(0.35))
            self._run(tf.paragraphs[0], spec["heading"], font=BODY_FONT, size=18, color=INK, bold=True)

            visual_y = top + Inches(0.75)
            visual_h = Inches(1.6)
            if kind == "relational":
                self._table_grid(s, x + Inches(0.5), visual_y, col_w - Inches(1.0), visual_h, rows=4, cols=3)
            else:
                self._document_icon(s, x + Inches(1.15), visual_y, col_w - Inches(2.3), visual_h)

            _, tf = self._textbox(s, x + Inches(0.3), visual_y + visual_h + Inches(0.2), col_w - Inches(0.6), Inches(0.9))
            tf.word_wrap = True
            self._run(tf.paragraphs[0], spec["desc"], font=BODY_FONT, size=13.5, color=GRAY)

            _, tf = self._textbox(s, x + Inches(0.3), top + Inches(3.6), col_w - Inches(0.6), Inches(0.4))
            self._run(tf.paragraphs[0], "EXAMPLES  ", font=MONO_FONT, size=10.5, color=BLUE, bold=True)
            self._run(tf.paragraphs[0], spec["examples"], font=BODY_FONT, size=12.5, color=INK)

        self._page_number(s)
        return s

    def entity_attribute_slide(self, eyebrow, title, subtitle, left_entity, right_entity, relationship, left_attrs, right_attrs):
        """Chen-notation ER diagram: [attrs] - Entity - <relationship diamond> - Entity - [attrs]."""
        s = self._slide()
        self._header(s, eyebrow, title, subtitle)

        cy = Inches(4.3)
        ent_w, ent_h = Inches(2.0), Inches(0.7)
        lx = Inches(2.3)
        rx = Inches(10.4)
        dw, dh = Inches(1.7), Inches(0.95)
        dcx = (lx + rx) / 2 + ent_w / 2

        self._entity_box(s, lx, cy - ent_h / 2, ent_w, ent_h, left_entity, fill=WHITE, border=BLUE)
        self._entity_box(s, rx, cy - ent_h / 2, ent_w, ent_h, right_entity, fill=WHITE, border=BLUE)
        self._connector(s, lx + ent_w, cy, dcx - dw / 2, cy, color=GRAY_MUTED, width=1.25)
        self._connector(s, dcx + dw / 2, cy, rx, cy, color=GRAY_MUTED, width=1.25)
        self._diamond(s, dcx, cy, dw, dh, relationship)

        oval_w, oval_h = Inches(1.3), Inches(0.55)
        stride = Inches(1.55)
        ay = cy - Inches(1.6)
        for i, attr in enumerate(left_attrs):
            ax = lx + ent_w / 2 + stride * (i - (len(left_attrs) - 1) / 2)
            self._oval(s, ax, ay, oval_w, oval_h, attr)
            self._connector(s, lx + ent_w / 2, cy - ent_h / 2, ax, ay + oval_h / 2, color=GRAY_MUTED, width=1.0)
        for i, attr in enumerate(right_attrs):
            ax = rx + ent_w / 2 + stride * (i - (len(right_attrs) - 1) / 2)
            self._oval(s, ax, ay, oval_w, oval_h, attr)
            self._connector(s, rx + ent_w / 2, cy - ent_h / 2, ax, ay + oval_h / 2, color=GRAY_MUTED, width=1.0)

        legend_y = Inches(6.35)
        items = [("Rectangle = Entity", BLUE), ("Oval = Attribute", CARD_BORDER), ("Diamond = Relationship", BLUE)]
        lx0 = Inches(3.0)
        for i, (label, color) in enumerate(items):
            self._caption(s, lx0 + Inches(3.0) * i, legend_y, Inches(2.9), label, size=12, color=GRAY)

        self._page_number(s)
        return s

    def cards_slide(self, eyebrow, title, subtitle, cards, panel_title=None, panel_lines=None):
        """cards: list of (heading, desc) stacked in the left column.
        panel_lines: list of strings (or (bold_lead, rest) tuples) in the right gray panel."""
        s = self._slide()
        self._header(s, eyebrow, title, subtitle)

        left_x = MARGIN
        card_w = Inches(5.7) if panel_lines else Inches(11.9)
        top = Inches(2.7)
        card_h = Inches(1.05)
        gap = Inches(0.22)
        for i, (heading, desc) in enumerate(cards):
            y = top + (card_h + gap) * i
            self._rounded_rect(s, left_x, y, card_w, card_h, fill=WHITE, line_color=CARD_BORDER)
            _, tf = self._textbox(s, left_x + Inches(0.3), y + Inches(0.15), card_w - Inches(0.6), Inches(0.4))
            self._run(tf.paragraphs[0], heading, font=BODY_FONT, size=18, color=INK, bold=True)
            _, tf = self._textbox(s, left_x + Inches(0.3), y + Inches(0.58), card_w - Inches(0.6), Inches(0.4))
            self._run(tf.paragraphs[0], desc, font=BODY_FONT, size=14, color=GRAY)

        if panel_lines:
            panel_x = Inches(6.63)
            panel_w = Inches(6.1)
            panel_h = (card_h + gap) * len(cards) - gap
            self._rounded_rect(s, panel_x, top, panel_w, panel_h, fill=PANEL_BG)
            _, tf = self._textbox(s, panel_x + Inches(0.35), top + Inches(0.3), panel_w - Inches(0.7), Inches(0.35))
            self._run(tf.paragraphs[0], (panel_title or "").upper(), font=MONO_FONT, size=13, color=BLUE, bold=True)

            _, tf = self._textbox(s, panel_x + Inches(0.35), top + Inches(0.8), panel_w - Inches(0.7), panel_h - Inches(1.1))
            tf.word_wrap = True
            first = True
            for line in panel_lines:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.space_after = Pt(14)
                self._run(p, line, font=BODY_FONT, size=15.5, color=INK)

        self._page_number(s)
        return s

    def relationship_diagram_slide(self, eyebrow, title, subtitle, diagrams):
        """diagrams: list of exactly 3 dicts, each:
            {kind: '1:1'|'1:N'|'N:N', left_label, right_label, caption}
        Draws a small box-and-line diagram per type, side by side."""
        s = self._slide()
        self._header(s, eyebrow, title, subtitle)

        col_w = Inches(3.9)
        gap = Inches(0.3)
        top = Inches(2.9)
        for i, d in enumerate(diagrams):
            x0 = MARGIN + (col_w + gap) * i
            cy = top + Inches(1.0)
            box_w, box_h = Inches(1.35), Inches(0.6)
            kind = d["kind"]

            if kind == "1:1":
                lx = x0 + Inches(0.15)
                rx = x0 + col_w - box_w - Inches(0.15)
                self._entity_box(s, lx, cy - box_h / 2, box_w, box_h, d["left_label"])
                self._entity_box(s, rx, cy - box_h / 2, box_w, box_h, d["right_label"])
                self._connector(s, lx + box_w, cy, rx, cy, color=BLUE, width=1.5)
                self._tag(s, lx + box_w + Inches(0.28), cy - Inches(0.3), "1")
                self._tag(s, rx - Inches(0.28), cy - Inches(0.3), "1")

            elif kind == "1:N":
                lx = x0 + Inches(0.15)
                self._entity_box(s, lx, cy - box_h / 2, box_w, box_h, d["left_label"])
                rx = x0 + col_w - Inches(1.15)
                small_w, small_h = Inches(1.0), Inches(0.42)
                ys = [cy - Inches(0.55), cy, cy + Inches(0.55)]
                for j, ry in enumerate(ys):
                    self._entity_box(s, rx, ry - small_h / 2, small_w, small_h, f'{d["right_label"]} {j + 1}', fill=PANEL_BG, border=CARD_BORDER)
                    self._connector(s, lx + box_w, cy, rx, ry, color=GRAY_MUTED, width=1.1)
                self._tag(s, lx + box_w + Inches(0.22), cy - Inches(0.3), "1")
                self._tag(s, rx - Inches(0.22), cy - Inches(0.3), "N")

            else:  # N:N
                lx = x0 + Inches(0.15)
                rx = x0 + col_w - Inches(1.5)
                small_w, small_h = Inches(1.35), Inches(0.42)
                left_ys = [cy - Inches(0.4), cy + Inches(0.4)]
                right_ys = [cy - Inches(0.4), cy + Inches(0.4)]
                for j, ly in enumerate(left_ys):
                    self._entity_box(s, lx, ly - small_h / 2, small_w, small_h, f'{d["left_label"]} {j + 1}', fill=PANEL_BG, border=CARD_BORDER)
                for k, ry in enumerate(right_ys):
                    self._entity_box(s, rx, ry - small_h / 2, small_w, small_h, f'{d["right_label"]} {k + 1}', fill=PANEL_BG, border=CARD_BORDER)
                for ly in left_ys:
                    for ry in right_ys:
                        self._connector(s, lx + small_w, ly, rx, ry, color=GRAY_MUTED, width=1.0)
                self._tag(s, lx + small_w + Inches(0.22), cy - Inches(0.62), "N")
                self._tag(s, rx - Inches(0.22), cy - Inches(0.62), "N")

            self._caption(s, x0, top + Inches(1.95), col_w, d["caption"], size=12.5, color=GRAY)
            _, tf = self._textbox(s, x0, top - Inches(0.35), col_w, Inches(0.3))
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._run(tf.paragraphs[0], kind, font=MONO_FONT, size=13, color=BLUE, bold=True)

        self._page_number(s)
        return s

    def erd_slide(self, eyebrow, title, subtitle, center, branches):
        """center: (name, columns_str). branches: list of (name, columns_str, cardinality) —
        cardinality like 'N' drawn on the branch end, '1' always drawn on the center end."""
        s = self._slide()
        self._header(s, eyebrow, title, subtitle)

        cw, ch = Inches(2.3), Inches(0.85)
        cx = Inches(1.1)
        cy = Inches(4.3)
        self._entity_box(s, cx, cy - ch / 2, cw, ch, center[0], sub=center[1], fill=PANEL_BG, border=BLUE)

        bw, bh = Inches(2.5), Inches(0.85)
        bx = Inches(8.0)
        n = len(branches)
        spacing = Inches(2.0)
        start_y = cy - spacing * (n - 1) / 2
        for i, (name, cols, card) in enumerate(branches):
            by = start_y + spacing * i
            self._entity_box(s, bx, by - bh / 2, bw, bh, name, sub=cols, fill=WHITE, border=CARD_BORDER)
            self._connector(s, cx + cw, cy, bx, by, color=GRAY_MUTED, width=1.25)
            near_center_y = cy + (by - cy) * 0.18
            self._tag(s, cx + cw + Inches(0.35), near_center_y, "1", diameter=0.3)
            self._tag(s, bx - Inches(0.35), by, card, diameter=0.3)

        self._page_number(s)
        return s

    def limit_offset_diagram_slide(self, eyebrow, title, subtitle, total_rows, offset, limit):
        s = self._slide()
        self._header(s, eyebrow, title, subtitle)

        row_w = Inches(0.85)
        row_h = Inches(0.55)
        gap = Inches(0.08)
        start_x = MARGIN + Inches(0.3)
        y = Inches(3.6)

        for i in range(total_rows):
            x = start_x + (row_w + gap) * i
            if i < offset:
                fill, border, txt_color = PANEL_BG, CARD_BORDER, GRAY_MUTED
            elif i < offset + limit:
                fill, border, txt_color = BLUE, BLUE, WHITE
            else:
                fill, border, txt_color = WHITE, CARD_BORDER, GRAY_MUTED
            self._rounded_rect(s, x, y, row_w, row_h, fill=fill, line_color=border, radius=0.15)
            _, tf = self._textbox(s, x, y, row_w, row_h)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._run(tf.paragraphs[0], f"row {i + 1}", font=MONO_FONT, size=10, color=txt_color, bold=(fill == BLUE))

        bracket_y = y - Inches(0.35)
        if offset > 0:
            bx0 = start_x
            bx1 = start_x + (row_w + gap) * offset - gap
            self._connector(s, bx0, bracket_y, bx1, bracket_y, color=GRAY_MUTED, width=1.25)
            self._caption(s, bx0, bracket_y - Inches(0.38), bx1 - bx0, f"OFFSET {offset} — skipped", size=12, color=GRAY)

        lx0 = start_x + (row_w + gap) * offset
        lx1 = start_x + (row_w + gap) * (offset + limit) - gap
        self._connector(s, lx0, bracket_y, lx1, bracket_y, color=BLUE, width=1.75)
        self._caption(s, lx0, bracket_y - Inches(0.38), lx1 - lx0, f"LIMIT {limit} — returned", size=12, color=BLUE, bold=True)

        self._page_number(s)
        return s

    def code_slide(self, eyebrow, title, subtitle, filename, code_lines, notes=None):
        """code_lines: list of raw SQL lines (auto-tokenized) OR (text, [(t,color),...]) pre-tokenized.
        notes: list of (heading, desc) shown to the right of the code block."""
        s = self._slide()
        self._header(s, eyebrow, title, subtitle)

        code_x = MARGIN
        code_w = Inches(7.3) if notes else Inches(11.9)
        code_y = Inches(2.75)
        header_h = Inches(0.5)
        code_h = Inches(3.9)

        self._rounded_rect(s, code_x, code_y, code_w, header_h + code_h, fill=CODE_BG, radius=0.06)

        for i, color in enumerate(TRAFFIC):
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, code_x + Inches(0.25 + i * 0.28), code_y + Inches(0.17), Inches(0.16), Inches(0.16))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            dot.shadow.inherit = False
        _, tf = self._textbox(s, code_x + Inches(1.1), code_y + Inches(0.13), code_w - Inches(1.3), Inches(0.3))
        self._run(tf.paragraphs[0], filename, font=MONO_FONT, size=12, color=RGBColor(0x9C, 0x9C, 0x9C))

        _, tf = self._textbox(s, code_x + Inches(0.35), code_y + header_h + Inches(0.25), code_w - Inches(0.6), code_h - Inches(0.4))
        tf.word_wrap = False
        first = True
        for line in code_lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(4)
            if isinstance(line, tuple):
                text, tokens = line
                for t, color in tokens:
                    self._run(p, t, font=MONO_FONT, size=14, color=color)
            elif line == "":
                self._run(p, " ", font=MONO_FONT, size=14, color=CODE_DEFAULT)
            else:
                for t, color in _tokenize_sql_line(line):
                    self._run(p, t, font=MONO_FONT, size=14, color=color)

        if notes:
            notes_x = Inches(8.6)
            notes_w = Inches(4.13)
            chars_per_line = 46  # rough estimate for 13.5pt Nunito Sans in a 4.13in column
            y = code_y
            for heading, desc in notes:
                _, tf = self._textbox(s, notes_x, y, notes_w, Inches(0.35))
                self._run(tf.paragraphs[0], heading, font=BODY_FONT, size=16, color=INK, bold=True)
                import math
                n_lines = max(1, math.ceil(len(desc) / chars_per_line))
                desc_h = Inches(0.24 * n_lines + 0.08)
                _, tf = self._textbox(s, notes_x, y + Inches(0.35), notes_w, desc_h)
                tf.word_wrap = True
                self._run(tf.paragraphs[0], desc, font=BODY_FONT, size=13.5, color=GRAY)
                y += Inches(0.35) + desc_h + Inches(0.22)

        self._page_number(s)
        return s

    def closing_slide(self, eyebrow, title, subtitle):
        s = self._slide()
        cy = Inches(3.2)
        _, tf = self._textbox(s, Inches(0), cy, SLIDE_W, Inches(0.4))
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        self._run(tf.paragraphs[0], eyebrow.upper(), font=MONO_FONT, size=15, color=BLUE, bold=True)

        _, tf = self._textbox(s, Inches(0.5), cy + Inches(0.5), Inches(12.333), Inches(1.1))
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        self._run(tf.paragraphs[0], title, font=TITLE_FONT, size=46, color=INK, bold=True)

        _, tf = self._textbox(s, Inches(1.5), cy + Inches(1.55), Inches(10.333), Inches(0.6))
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        self._run(tf.paragraphs[0], subtitle, font=BODY_FONT, size=18, color=GRAY)
        return s

    # -- finalize ---------------------------------------------------------

    def save(self, path):
        total = self._n
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "{TOTAL}" in run.text:
                            run.text = run.text.replace("{TOTAL}", str(total))
        self.prs.save(path)
        return path
